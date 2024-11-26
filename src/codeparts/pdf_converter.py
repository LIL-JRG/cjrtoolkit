import os
import io
import tkinter as tk
from tkinter import filedialog
from PyPDF2 import PdfReader
from docx import Document
from docx.shared import Inches
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches as PptxInches
from PIL import Image
import pdfplumber
import tempfile
import camelot
import fitz
import pandas as pd
from system.config import RESULT_FOLDER, PDF_CONVERTER_CONFIG
import win32com.client
import pythoncom
import comtypes.client
import pdf2docx
import logging
import comtypes.gen
from pdf2image import convert_from_path

comtypes_logger = logging.getLogger('comtypes')
comtypes_logger.setLevel(logging.WARNING)

logging.basicConfig(level=logging.INFO)

class PDFConverter:
    @staticmethod
    def select_file():
        root = tk.Tk()
        root.withdraw()
        file_path = filedialog.askopenfilename(
            filetypes=PDF_CONVERTER_CONFIG["supported_formats"]["input"]
        )
        return file_path

    @staticmethod
    def create_output_folder(input_file):
        pdf_converter_folder = os.path.join(RESULT_FOLDER, "PDF Converter")
        os.makedirs(pdf_converter_folder, exist_ok=True)
        
        file_name = os.path.splitext(os.path.basename(input_file))[0]
        output_folder = os.path.join(pdf_converter_folder, file_name)
        os.makedirs(output_folder, exist_ok=True)
        
        return output_folder

    @staticmethod
    async def pdf_to_word(pdf_path: str) -> str:
        output_folder = PDFConverter.create_output_folder(pdf_path)
        output_path = os.path.join(output_folder, f"{os.path.splitext(os.path.basename(pdf_path))[0]}.docx")
        try:
            converter = pdf2docx.Converter(pdf_path)
            converter.convert(output_path)
            converter.close()
            return f"Conversión exitosa. Archivo guardado en: {output_path}"
        except Exception as e:
            return f"Error en la conversión a Word: {str(e)}"

    @staticmethod
    async def pdf_to_excel(pdf_path: str) -> str:
        output_folder = PDFConverter.create_output_folder(pdf_path)
        output_path = os.path.join(output_folder, f"{os.path.splitext(os.path.basename(pdf_path))[0]}.xlsx")
        try:
            tables = camelot.read_pdf(
                pdf_path,
                flavor=PDF_CONVERTER_CONFIG["table_settings"]["flavor"],
                line_scale=PDF_CONVERTER_CONFIG["table_settings"]["line_scale"]
            )
            
            with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                for i, table in enumerate(tables):
                    df = table.df
                    df.to_excel(writer, sheet_name=f'Table_{i+1}', index=False)
            
            return f"Conversión exitosa. Archivo guardado en: {output_path}"
        except Exception as e:
            return f"Error en la conversión a Excel: {str(e)}"

    @staticmethod
    async def pdf_to_powerpoint(pdf_path: str) -> str:
        output_folder = PDFConverter.create_output_folder(pdf_path)
        output_path = os.path.join(output_folder, f"{os.path.splitext(os.path.basename(pdf_path))[0]}.pptx")
        try:
            pdf = fitz.open(pdf_path)
            prs = Presentation()
        
            for page_num in range(len(pdf)):
                page = pdf[page_num]
                slide = prs.slides.add_slide(prs.slide_layouts[6])
            
                text = page.get_text()
                txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.5), 
                                                 PptxInches(9), PptxInches(5))
                tf = txBox.text_frame
                tf.text = text
            
                images = page.get_images(full=True)
                for img_index, img in enumerate(images):
                    xref = img[0]
                    base_image = pdf.extract_image(xref)
                    image_bytes = base_image["image"]
                
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
                        temp_file.write(image_bytes)
                        temp_file_path = temp_file.name
                
                    slide.shapes.add_picture(temp_file_path, PptxInches(1), PptxInches(3))
                
                    os.unlink(temp_file_path)
        
            prs.save(output_path)
            return f"Conversión exitosa. Archivo guardado en: {output_path}"
        except Exception as e:
            return f"Error en la conversión a PowerPoint: {str(e)}"

    @staticmethod
    async def extract_images_from_pdf(pdf_path: str) -> str:
        output_folder = PDFConverter.create_output_folder(pdf_path)
        try:
            images = convert_from_path(
                pdf_path,
                dpi=PDF_CONVERTER_CONFIG["image_settings"]["dpi"]
            )
            
            image_paths = []
            for i, image in enumerate(images):
                image_path = os.path.join(
                    output_folder,
                    f'page_{i+1}.{PDF_CONVERTER_CONFIG["image_settings"]["format"].lower()}'
                )
                image.save(image_path, PDF_CONVERTER_CONFIG["image_settings"]["format"])
                image_paths.append(image_path)
            
            return f"Extracción de imágenes exitosa. Imágenes guardadas en: {output_folder}"
        except Exception as e:
            return f"Error en la extracción de imágenes: {str(e)}"

    @staticmethod
    async def office_to_pdf(input_path: str) -> str:
        output_folder = PDFConverter.create_output_folder(input_path)
        output_path = os.path.join(output_folder, f"{os.path.splitext(os.path.basename(input_path))[0]}.pdf")
        try:
            pythoncom.CoInitialize()
            if input_path.endswith('.docx'):
                word = win32com.client.Dispatch('Word.Application')
                doc = word.Documents.Open(input_path)
                doc.SaveAs(output_path, FileFormat=17)
                doc.Close()
                word.Quit()
            elif input_path.endswith('.xlsx'):
                excel = win32com.client.Dispatch('Excel.Application')
                wb = excel.Workbooks.Open(input_path)
                wb.ExportAsFixedFormat(0, output_path)
                wb.Close()
                excel.Quit()
            elif input_path.endswith('.pptx'):
                powerpoint = win32com.client.Dispatch('Powerpoint.Application')
                presentation = powerpoint.Presentations.Open(input_path)
                presentation.SaveAs(output_path, 32)
                presentation.Close()
                powerpoint.Quit()
            else:
                return "Formato de archivo no soportado"
            return f"Conversión exitosa. Archivo guardado en: {output_path}"
        except Exception as e:
            return f"Error en la conversión a PDF: {str(e)}"
        finally:
            pythoncom.CoUninitialize()